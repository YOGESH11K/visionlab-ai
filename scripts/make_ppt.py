#!/usr/bin/env python
"""Generate VisionLab AI (Empire) presentation deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- palette ----------
BG      = RGBColor(10, 12, 18)     # deep space blue-black
PANEL   = RGBColor(18, 22, 32)     # glass panel
PANEL2  = RGBColor(24, 30, 44)
ACCENT  = RGBColor(0, 200, 180)    # cyan-teal
ACCENT2 = RGBColor(120, 190, 255)  # light blue
ORANGE  = RGBColor(255, 150, 70)
RED     = RGBColor(255, 90, 90)
GREEN   = RGBColor(90, 220, 120)
YELLOW  = RGBColor(240, 210, 90)
INK     = RGBColor(235, 240, 248)
DIM     = RGBColor(150, 160, 180)
DARKER  = RGBColor(8, 10, 15)

FONT = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = prs.slide_width, prs.slide_height

def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=PANEL, line=None, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, italic, mono)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, mono) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = MONO if mono else FONT
    return tb

def header(slide, kicker, title, idx=None, total=None):
    # top rule
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background(); ln.shadow.inherit = False
    add_text(slide, 0.55, 0.42, 9.0, 0.32,
             [[(kicker.upper(), 11, ACCENT, True, False, False)]])
    add_text(slide, 0.55, 0.72, 11.5, 0.7,
             [[(title, 30, INK, True, False, False)]])
    if idx is not None:
        add_text(slide, 12.0, 0.55, 1.0, 0.4,
                 [[(f"{idx:02d}", 13, DIM, True, False, True)]], align=PP_ALIGN.RIGHT)
    # hairline under header
    ln2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.5), Inches(12.25), Pt(1))
    ln2.fill.solid(); ln2.fill.fore_color.rgb = PANEL2; ln2.line.fill.background(); ln2.shadow.inherit = False

def bullet_items(items, accent=ACCENT, size=15, gap=10):
    """items: list of (bold_part, rest)"""
    runs = []
    for bold, rest in items:
        runs.append([("▸  ", size, accent, True, False, False),
                     (bold, size, INK, True, False, False),
                     (rest, size, DIM, False, False, False)])
        runs.append([("", 4, INK, False, False, False)]) if False else None
    return runs

# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
# grid decoration
for i in range(0, 14):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i*1.0), Inches(0), Pt(1), SH)
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(20, 26, 38); ln.line.fill.background(); ln.shadow.inherit = False
for j in range(0, 8):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(j*1.1), SW, Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(20, 26, 38); ln.line.fill.background(); ln.shadow.inherit = False
# glow line
glow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(2.35), Inches(2.2), Pt(3))
glow.fill.solid(); glow.fill.fore_color.rgb = ACCENT; glow.line.fill.background(); glow.shadow.inherit = False

add_text(s, 0.55, 1.55, 12.0, 0.4, [[("VISIONLAB AI  ·  ENGINEERING LAB", 14, ACCENT, True, False, True)]])
add_text(s, 0.55, 2.6, 12.2, 2.0,
         [[("AI-Powered Electronics Laboratory", 54, INK, True, False, False)],
          [("Computer Vision  ·  Hand Tracking  ·  Arduino/ESP32  ·  Sensors  ·  Learning", 20, DIM, False, False, False)]])
add_text(s, 0.55, 4.55, 11.0, 0.5,
         [[("Where the camera understands the physical world, and you control real hardware — naturally.", 17, DIM, False, True, False)]])

# three chips
chips = [("👁  Vision", "MediaPipe hand tracking"), ("🔌  Hardware", "Arduino · ESP32 · Virtual"), ("🤖  AI", "Assistant + Code Generator")]
cx = 0.55
for label, sub in chips:
    add_rect(s, cx, 5.35, 3.85, 0.95, PANEL, line=PANEL2)
    add_text(s, cx+0.25, 5.52, 3.5, 0.7,
             [[(label + "  ", 16, ACCENT, True, False, False)],
              [(sub, 12, DIM, False, False, False)]])
    cx += 4.1

add_text(s, 0.55, 6.75, 12.0, 0.4,
         [[("Live:  visionlab-ai.vercel.app   ·   github.com/YOGESH11K/visionlab-ai", 12, DIM, False, False, True)]])

# =====================================================================
# SLIDE 2 — WHY WE NEED THIS
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "The problem", "Why we need this project", 2, 15)
add_text(s, 0.55, 1.75, 12.2, 0.6,
         [[("Learning electronics & robotics today is fragmented — theory, code, wiring and hardware never live in one place.", 17, INK, False, False, False)]])

problems = [
    ("📚", "Knowledge is scattered", "Components, wiring diagrams and code live in different tutorials, PDFs and videos.", ACCENT2),
    ("🧩", "Theory vs. reality gap", "Students learn circuits in theory but never see the real hardware respond.", ORANGE),
    ("🐌", "Slow feedback loops", "Every change means: edit code → upload → run → debug. Painful for beginners.", YELLOW),
    ("💸", "Hardware is expensive", "Not everyone owns an Arduino, sensors and a working setup.", RED),
    ("🚧", "No unified workspace", "No single tool combines camera, gestures, circuits, code and hardware control.", DIM),
]
y = 2.5
for icon, title, desc, col in problems:
    add_rect(s, 0.55, y, 12.25, 0.78, PANEL, line=PANEL2)
    add_text(s, 0.85, y+0.14, 0.6, 0.5, [[(icon, 20, col, False, False, False)]])
    add_text(s, 1.55, y+0.13, 3.4, 0.5, [[(title, 16, col, True, False, False)]])
    add_text(s, 5.0, y+0.14, 7.6, 0.55, [[(desc, 14, DIM, False, False, False)]])
    y += 0.9
add_text(s, 0.55, 7.0, 12.2, 0.4, [[("VisionLab AI solves all five: one unified, gesture-controlled, AI-assisted electronics laboratory.", 15, ACCENT, True, False, False)]])

# =====================================================================
# SLIDE 3 — WHAT IT IS / SOLUTION
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "The solution", "One professional electronics laboratory", 3, 15)
add_text(s, 0.55, 1.75, 12.2, 0.6,
         [[("VisionLab AI is a production-grade platform where your camera becomes the interface to real and simulated hardware.", 17, INK, False, False, False)]])

sol = [
    ("CAMERA", "computer vision understands the physical world", ACCENT),
    ("GESTURES", "control LEDs, motors, servos with hand gestures", ACCENT2),
    ("HARDWARE", "Arduino / ESP32 control, or a full Virtual Arduino", ORANGE),
    ("SENSORS", "live distance, temperature, humidity, light & motion", GREEN),
    ("AI", "ask questions, generate Arduino code instantly", YELLOW),
    ("CIRCUITS", "visual builder with smart validation", DIM),
]
col = 0; row = 0
for k, v, c in sol:
    x = 0.55 + col*4.15; y = 2.6 + row*1.5
    add_rect(s, x, y, 3.9, 1.25, PANEL, line=PANEL2)
    add_rect(s, x, y, 0.09, 1.25, c)
    add_text(s, x+0.3, y+0.18, 3.4, 0.4, [[(k, 15, c, True, False, True)]])
    add_text(s, x+0.3, y+0.6, 3.5, 0.55, [[(v, 12.5, DIM, False, False, False)]])
    col += 1
    if col > 2:
        col = 0; row += 1
add_text(s, 0.55, 6.6, 12.2, 0.4, [[("Result: a futuristic engineering-lab experience on desktop, laptop or tablet.", 15, ACCENT, True, False, False)]])

# =====================================================================
# SLIDE 4 — HOW IT WORKS (PIPELINE)
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Architecture", "How it works — the full pipeline", 4, 15)
steps = [
    ("CAMERA", "frame capture\n(sim fallback)"),
    ("VISION", "MediaPipe hand\ntracking · 21 points"),
    ("UNDERSTAND", "finger count\ngesture engine"),
    ("UI", "live dashboard\noverlays + events"),
    ("COMMAND", "gesture→action\nmapping (stable)"),
    ("HARDWARE", "Arduino / ESP32\nor Virtual board"),
]
x = 0.55
for i, (t, sub) in enumerate(steps):
    add_rect(s, x, 2.1, 1.95, 1.5, PANEL, line=PANEL2)
    add_rect(s, x, 2.1, 1.95, 0.12, ACCENT if i < 4 else ORANGE)
    add_text(s, x+0.15, 2.35, 1.7, 0.4, [[(t, 13.5, ACCENT if i<4 else ORANGE, True, False, True)]])
    add_text(s, x+0.15, 2.8, 1.7, 0.7, [[(l, 11.5, DIM, False, False, False)] for l in sub.split("\n")])
    if i < len(steps)-1:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x+1.98), Inches(2.62), Inches(0.22), Inches(0.45))
        ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor(60, 70, 90); ar.line.fill.background(); ar.shadow.inherit = False
    x += 2.17
# example card
add_rect(s, 0.55, 4.15, 12.25, 1.35, PANEL2, line=PANEL2)
add_text(s, 0.85, 4.3, 11.8, 0.35, [[("LIVE EXAMPLE — “Show 3 fingers”", 12.5, DIM, True, False, True)]])
add_text(s, 0.85, 4.75, 11.9, 0.7,
         [[("Show 3 fingers  →  ", 14, INK, False, False, False), ("MediaPipe detects hand", 14, ACCENT, True, False, False),
           ("  →  ", 14, DIM, False, False, False), ("Finger counter = 3", 14, ACCENT2, True, False, False),
           ("  →  ", 14, DIM, False, False, False), ("Gesture engine", 14, YELLOW, True, False, False),
           ("  →  ", 14, DIM, False, False, False), ("LED_3_ON", 14, ORANGE, True, False, True)]])

add_text(s, 0.55, 6.1, 12.2, 0.9,
         [[("Same pipeline powers component recognition:", 14, INK, False, False, False)],
          [("Camera sees HC-SR04  →  vision identifies it  →  info panel (pins, wiring, Arduino example)  →  connect sensor  →  live distance appears.", 13.5, DIM, False, False, False)]])

# =====================================================================
# SLIDE 5 — TECH STACK
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Technology", "Built on a simple, reliable, modular stack", 5, 15)
stack = [
    ("FRONTEND", "React · TypeScript · Vite · Tailwind CSS", "Dark engineering-lab UI, 12 workspaces, live WebSockets", ACCENT),
    ("VISION", "Python · OpenCV · MediaPipe", "21-point hand tracking, gestures, bounding boxes, AR overlay", ACCENT2),
    ("BACKEND", "Python · FastAPI · WebSocket · REST", "Vision/hardware/sensor services, event bus, SQLite (PostgreSQL-ready)", ORANGE),
    ("HARDWARE", "Arduino · ESP32 · pyserial", "Structured command protocol with IDs, Virtual Arduino simulator", GREEN),
    ("DATA", "SQLite via SQLAlchemy", "Mappings, sensors, events, projects, scores — Postgres-swappable", YELLOW),
    ("AI", "Verified knowledge engine + optional LLM", "No hallucinated specs — unknown answers say “unknown”", DIM),
]
y = 1.75
for name, stackline, desc, col in stack:
    add_rect(s, 0.55, y, 12.25, 0.82, PANEL, line=PANEL2)
    add_rect(s, 0.55, y, 0.09, 0.82, col)
    add_text(s, 0.85, y+0.1, 3.6, 0.6, [[(name, 14, col, True, False, True)], [(stackline, 11, DIM, False, False, True)]])
    add_text(s, 4.7, y+0.2, 8.0, 0.5, [[(desc, 13, INK, False, False, False)]])
    y += 0.9

# =====================================================================
# SLIDE 6 — GESTURE CONTROL
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Hand tracking & gesture control", "Control hardware with your hands", 6, 15)
add_text(s, 0.55, 1.75, 12.2, 0.6,
         [[("Real MediaPipe hand tracking: 21 landmarks, left/right detection, finger states, palm bounding box, confidence score.", 15, DIM, False, False, False)]])

# default mapping
add_rect(s, 0.55, 2.5, 6.0, 3.4, PANEL, line=PANEL2)
add_text(s, 0.85, 2.68, 5.5, 0.35, [[("DEFAULT GESTURE → ACTION", 12.5, ACCENT, True, False, True)]])
maps = [("0 fingers", "ALL LEDs OFF"), ("1 finger", "LED 1 ON"), ("2 fingers", "LED 2 ON"),
        ("3 fingers", "LED 3 ON"), ("4 fingers", "LED 4 ON"), ("5 fingers", "ALL LEDs ON")]
yy = 3.15
for g, a in maps:
    add_text(s, 1.05, yy, 2.6, 0.35, [[(g, 14, INK, True, False, False)]])
    add_text(s, 3.6, yy, 2.8, 0.35, [[("→  " + a, 14, ACCENT2, True, False, False)]])
    yy += 0.44
add_text(s, 0.85, 5.85, 5.5, 0.5, [[("Fully configurable — no hard-coded mappings.", 12, DIM, False, True, False)]])

# advanced gestures
add_rect(s, 6.85, 2.5, 5.95, 3.4, PANEL, line=PANEL2)
add_text(s, 7.15, 2.68, 5.5, 0.35, [[("ADVANCED GESTURES", 12.5, ACCENT2, True, False, True)]])
adv = [("OPEN PALM", "all on"), ("FIST", "emergency off"), ("THUMB UP", "confirm / boost"),
       ("PEACE", "change control mode"), ("PINCH", "brightness (PWM)"), ("SWIPE LEFT / RIGHT", "previous / next mode")]
yy = 3.15
for g, a in adv:
    add_text(s, 7.35, yy, 3.4, 0.35, [[(g, 13, INK, True, False, True)]])
    add_text(s, 10.7, yy, 2.0, 0.35, [[(a, 12.5, DIM, False, False, False)]])
    yy += 0.44
add_text(s, 7.15, 5.85, 5.5, 0.5,
         [[("Anti-spam engine: debounce + confidence threshold + temporal smoothing + cooldown — stable gestures only.", 12, DIM, False, True, False)]])

# =====================================================================
# SLIDE 7 — HARDWARE CONTROL
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Hardware lab", "Arduino, ESP32 & Virtual Arduino", 7, 15)
add_text(s, 0.55, 1.75, 12.2, 0.6,
         [[("A full hardware manager: auto-detect serial ports, connect/disconnect, board selection, baud rate, serial monitor, latency & error handling.", 15, DIM, False, False, False)]])

hw = [
    ("BOARDS", "Arduino Uno · Nano · Mega · ESP32 (extensible)", ACCENT),
    ("COMMANDS", "LED1_ON · LED3_OFF · LED1_PWM:120 · SERVO:90 · SENSOR", ACCENT2),
    ("PROTOCOL", "PC: COMMAND LED3_ON ID=1042  →  Arduino: OK ID=1042", ORANGE),
    ("OUTPUTS", "LEDs · PWM · servo · buzzer · relay · motor", YELLOW),
    ("INPUTS", "digital · analog · distance · temp · humidity · light · motion", GREEN),
    ("FIRMWARE", "Clean .ino firmware for Uno / ESP32 included in the repo", DIM),
]
x = 0.55; y = 2.6
for t, v, col in hw:
    add_rect(s, x, y, 6.0, 0.95, PANEL, line=PANEL2)
    add_rect(s, x, y, 0.09, 0.95, col)
    add_text(s, x+0.28, y+0.1, 2.1, 0.7, [[(t, 12.5, col, True, False, True)]])
    add_text(s, x+2.4, y+0.13, 3.5, 0.75, [[(v, 11.5, DIM, False, False, True if ":" in v or "ID=" in v else False)]])
    if x < 6.8:
        x = 6.85
    else:
        x = 0.55; y += 1.05

# Virtual Arduino
add_rect(s, 0.55, 5.55, 12.25, 1.35, PANEL2, line=ACCENT)
add_text(s, 0.85, 5.7, 11.7, 0.35, [[("NO HARDWARE? NO PROBLEM — VIRTUAL ARDUINO", 13, ACCENT, True, False, True)]])
add_text(s, 0.85, 6.15, 11.8, 0.7,
         [[("The built-in simulator behaves exactly like a real board — LEDs, servo, buzzer, inputs, sensors. Same command protocol, same UI. ", 13.5, INK, False, False, False),
           ("Demo-ready anywhere.", 13.5, GREEN, True, False, False)]])

# =====================================================================
# SLIDE 8 — SENSOR MONITORING
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Real-time data", "Sensor monitoring with professional charts", 8, 15)
add_text(s, 0.55, 1.75, 12.2, 0.6,
         [[("Live sensor values stream over WebSockets into professional engineering charts — current, min, max, average, trend and timestamp.", 15, DIM, False, False, False)]])

sens = [
    ("DHT11 / DHT22", "Temperature · Humidity", ACCENT),
    ("HC-SR04", "Ultrasonic distance", ACCENT2),
    ("LDR", "Light level", ORANGE),
    ("PIR", "Motion detection", YELLOW),
    ("POTENTIOMETER", "Analog value", GREEN),
]
x = 0.55
for name, desc, col in sens:
    add_rect(s, x, 2.5, 2.35, 1.25, PANEL, line=PANEL2)
    add_rect(s, x, 2.5, 2.35, 0.09, col)
    add_text(s, x+0.15, 2.68, 2.1, 0.4, [[(name, 12.5, col, True, False, True)]])
    add_text(s, x+0.15, 3.1, 2.1, 0.55, [[(desc, 11, DIM, False, False, False)]])
    x += 2.48

feats = [
    ("📈", "Charts", "Line · area · gauge · status indicators"),
    ("🎛️", "Controls", "Zoom · time range · pause · resume · clear"),
    ("📤", "Export", "CSV & JSON download, stored history"),
    ("⏱️", "Ranges", "last minute · 5 min · hour · today · custom"),
    ("📊", "Stats", "min · max · average · trend per channel"),
    ("🔔", "Live", "WebSocket stream — updates in real time"),
]
x = 0.55; y = 4.0
for icon, t, d in feats:
    add_rect(s, x, y, 3.95, 0.95, PANEL, line=PANEL2)
    add_text(s, x+0.2, y+0.1, 0.5, 0.4, [[(icon, 16, ACCENT, False, False, False)]])
    add_text(s, x+0.7, y+0.14, 3.1, 0.7, [[(t + "  ", 13.5, INK, True, False, False)], [(d, 11, DIM, False, False, False)]])
    x += 4.12
    if x > 12.4:
        x = 0.55; y += 1.05

# =====================================================================
# SLIDE 9 — COMPONENT RECOGNITION + AI
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Component recognition + AI assistant", "Point the camera · ask anything", 9, 15)

# left: component scan
add_rect(s, 0.55, 1.75, 6.0, 5.0, PANEL, line=PANEL2)
add_text(s, 0.85, 1.92, 5.5, 0.35, [[("COMPONENT SCANNER", 13, ACCENT, True, False, True)]])
add_text(s, 0.85, 2.35, 5.5, 1.2,
         [[("Point the camera at a component and the system suggests what it sees — with an honest confidence score. If unsure it shows “Possible match” and asks you to improve lighting / angle.", 12.5, DIM, False, False, False)]])
add_text(s, 0.85, 3.65, 5.5, 0.35, [[("SUPPORTED COMPONENTS (24+)", 12, ACCENT2, True, False, True)]])
comps = "Arduino Uno/Nano/Mega · ESP32 · LED · Resistor · Capacitor · Potentiometer · Push Button · LDR · PIR · HC-SR04 · DHT11 · DHT22 · IR Sensor · Servo · DC Motor · Buzzer · Relay · OLED · LCD"
add_text(s, 0.85, 4.05, 5.5, 1.3, [[(comps, 12, INK, False, False, False)]])
add_text(s, 0.85, 5.5, 5.5, 1.0,
         [[("Knowledge database stores pins, voltage, current, wiring, Arduino examples, ESP32 notes, mistakes & safety — displayed on detection.", 12, DIM, False, True, False)]])

# right: AI
add_rect(s, 6.85, 1.75, 5.95, 5.0, PANEL, line=PANEL2)
add_text(s, 7.15, 1.92, 5.5, 0.35, [[("AI ASSISTANT", 13, ACCENT2, True, False, True)]])
add_text(s, 7.15, 2.35, 5.5, 1.3,
         [[("Ask: “What is this?”, “How do I connect it?”, “Why isn’t it working?”, “Explain it to a beginner.”  Answers come from verified component data — it never hallucinates specs.", 12.5, DIM, False, False, False)]])
add_text(s, 7.15, 3.65, 5.5, 0.35, [[("CODE GENERATOR", 12, ORANGE, True, False, True)]])
add_text(s, 7.15, 4.05, 5.5, 1.0,
         [[("“Turn LED on below 10 cm.”  →  full Arduino sketch with wiring, pins, explanation and expected behaviour. Copy, save, validate, edit, download.", 12.5, DIM, False, False, False)]])
add_text(s, 7.15, 5.15, 5.5, 0.35, [[("SAFETY", 12, RED, True, False, True)]])
add_text(s, 7.15, 5.5, 5.5, 1.0,
         [[("No arbitrary code execution · no secret exposure · upload to hardware requires explicit confirmation.", 12, INK, False, False, False)]])

# =====================================================================
# SLIDE 10 — CIRCUIT BUILDER + LEARNING
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Design & learn", "Circuit builder + learning lab", 10, 15)

# circuit
add_rect(s, 0.55, 1.75, 6.0, 5.0, PANEL, line=PANEL2)
add_text(s, 0.85, 1.92, 5.5, 0.35, [[("CIRCUIT BUILDER", 13, ACCENT, True, False, True)]])
add_text(s, 0.85, 2.35, 5.5, 1.0,
         [[("Drag, drop, move, rotate and connect components — Arduino, ESP32, LED, sensors, servo, motor, buzzer, LCD, OLED. Pin labels shown on every part.", 12.5, DIM, False, False, False)]])
add_text(s, 0.85, 3.5, 5.5, 0.35, [[("SMART VALIDATION", 12, ACCENT2, True, False, True)]])
vv = [("🟢  VALID", "correct connections"), ("🟡  WARNING", "missing GND / power, voltage mismatch"), ("🔴  INVALID", "wrong pins, duplicates, short risk")]
yy = 3.9
for c, d in vv:
    add_text(s, 1.05, yy, 5.2, 0.35, [[(c, 12.5, INK, True, False, False), ("   " + d, 12, DIM, False, False, False)]])
    yy += 0.45
add_text(s, 0.85, 5.45, 5.5, 1.1,
         [[("Experimental camera analyzer compares expected vs observed circuit — marked MATCH / WARNING / POSSIBLE MISMATCH.", 12, DIM, False, True, False)]])

# learning
add_rect(s, 6.85, 1.75, 5.95, 5.0, PANEL, line=PANEL2)
add_text(s, 7.15, 1.92, 5.5, 0.35, [[("LEARNING LAB", 13, ACCENT2, True, False, True)]])
add_text(s, 7.15, 2.35, 5.5, 1.0,
         [[("Electronics learning platform with difficulty levels — BEGINNER · INTERMEDIATE · ADVANCED — for every component.", 12.5, DIM, False, False, False)]])
qz = ["Component quiz", "Pin quiz", "Circuit quiz", "Arduino coding quiz", "Scores & progress tracking", "Project suggestions"]
yy = 3.4
for q in qz:
    add_text(s, 7.35, yy, 5.3, 0.35, [[("▸  ", 13, ACCENT2, True, False, False), (q, 13, INK, False, False, False)]])
    yy += 0.42
add_text(s, 7.15, 5.1, 5.5, 1.1,
         [[("Pick your components (Arduino, LDR, LED, Servo…) and get suggested projects with concept, wiring, code and upgrades.", 12, DIM, False, True, False)]])

# =====================================================================
# SLIDE 11 — WORKSPACES / FEATURES MAP
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Product", "12 professional workspaces", 11, 15)
ws = ["Dashboard", "Vision Lab", "Gesture Control", "Component Scanner", "Sensor Monitor",
      "Circuit Builder", "Circuit Analyzer", "Arduino / ESP32", "AI Assistant", "Code Generator",
      "Projects", "Learning Lab", "History", "Settings"]
x = 0.55; y = 2.1
for i, w in enumerate(ws):
    add_rect(s, x, y, 2.9, 0.62, PANEL, line=PANEL2)
    add_text(s, x+0.2, y+0.12, 2.6, 0.4, [[(w, 12.5, INK if i < 12 else DIM, True, False, False)]])
    x += 3.1
    if x > 12.4:
        x = 0.55; y += 0.74
add_text(s, 0.55, 4.3, 12.2, 0.5,
         [[("Plus a real-time event console (time · source · event · command · status) and a diagnostics panel (FPS, latency, CPU, memory, WS status).", 14, DIM, False, False, False)]])
add_text(s, 0.55, 4.95, 12.2, 0.5,
         [[("Project management: save name, description, components, pin mappings, gestures, code, circuit, notes — 5 example projects pre-loaded.", 14, DIM, False, False, False)]])

# demo strip
add_rect(s, 0.55, 5.7, 12.25, 1.1, PANEL2, line=ACCENT)
add_text(s, 0.85, 5.85, 11.7, 0.35, [[("60-SECOND DEMO", 12.5, ACCENT, True, False, True)]])
add_text(s, 0.85, 6.25, 11.8, 0.5,
         [[("Show 1 finger → LED 1 ON · Show 3 fingers → LED 3 ON · Fist → all LEDs OFF · Show HC-SR04 → info + wiring · Ask AI → code · Build → simulate. ", 13, INK, False, False, False)]])

# =====================================================================
# SLIDE 12 — WHAT YOU CAN DO (USE CASES)
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Capabilities", "What you can do with VisionLab AI", 12, 15)
cases = [
    ("🖐", "Gesture-controlled LEDs", "Wave at the camera to switch lights, set brightness (PWM), sweep a servo."),
    ("📡", "Smart distance alarm", "HC-SR04 + AI: “turn the LED on below 10 cm” → code generated in seconds."),
    ("🌡️", "Temperature monitoring", "DHT11/DHT22 live charts, history and CSV export for reports."),
    ("🌙", "Automatic night light", "LDR brightness sensing with relay/motor control, all from the browser."),
    ("🚨", "Motion detection system", "PIR sensor events streamed live with console log and alarms."),
    ("🎓", "Learn electronics", "Quizzes, difficulty levels, project suggestions — from beginner to advanced."),
    ("🛠️", "Design circuits", "Visual builder + instant validation before you touch real hardware."),
    ("🧪", "Teach & demo anywhere", "Virtual Arduino means it runs with zero hardware in classrooms or demos."),
]
x = 0.55; y = 1.85
for icon, t, d in cases:
    add_rect(s, x, y, 5.95, 1.12, PANEL, line=PANEL2)
    add_text(s, x+0.2, y+0.08, 0.55, 0.4, [[(icon, 18, ACCENT, False, False, False)]])
    add_text(s, x+0.8, y+0.1, 5.0, 0.35, [[(t, 14, INK, True, False, False)]])
    add_text(s, x+0.8, y+0.48, 5.05, 0.55, [[(d, 11.5, DIM, False, False, False)]])
    if x < 6.8:
        x = 6.85
    else:
        x = 0.55; y += 1.24

# =====================================================================
# SLIDE 13 — BENEFITS / IMPACT
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Impact", "Why it matters", 13, 15)
benefits = [
    ("🎯", "One unified platform", "Vision + gestures + hardware + circuits + AI + learning — not a collection of separate demos.", ACCENT),
    ("💻", "Works without hardware", "Virtual Arduino + simulation camera make it fully demonstrable and affordable.", ACCENT2),
    ("⚡", "Instant feedback", "Gesture → LED responds in milliseconds; mistakes are caught immediately.", ORANGE),
    ("🧠", "Smarter learning", "AI explains like a teacher, quizzes measure progress, projects inspire next steps.", GREEN),
    ("🔒", "Safe by design", "No arbitrary code execution, no secret exposure, hardware uploads need confirmation.", YELLOW),
    ("🌍", "Anywhere, any device", "Runs in the browser on desktop, laptop and tablet — deployed on Vercel + Render.", DIM),
]
y = 1.85
for icon, t, d, col in benefits:
    add_rect(s, 0.55, y, 12.25, 0.85, PANEL, line=PANEL2)
    add_rect(s, 0.55, y, 0.09, 0.85, col)
    add_text(s, 0.85, y+0.1, 0.6, 0.5, [[(icon, 18, col, False, False, False)]])
    add_text(s, 1.55, y+0.09, 4.0, 0.5, [[(t, 15, col, True, False, False)]])
    add_text(s, 5.7, y+0.14, 7.0, 0.6, [[(d, 12.5, DIM, False, False, False)]])
    y += 0.93

# =====================================================================
# SLIDE 14 — DEPLOYMENT
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
header(s, "Live deployment", "Production-ready, already running", 14, 15)
add_rect(s, 0.55, 1.8, 6.0, 1.7, PANEL, line=ACCENT)
add_text(s, 0.85, 1.98, 5.4, 0.35, [[("🌐 FRONTEND — VERCEL", 13, ACCENT, True, False, True)]])
add_text(s, 0.85, 2.4, 5.4, 0.5, [[("https://visionlab-ai.vercel.app", 16, INK, True, False, True)]])
add_text(s, 0.85, 2.95, 5.4, 0.4, [[("React SPA · live video · WebSocket-ready", 12, DIM, False, False, False)]])

add_rect(s, 6.85, 1.8, 5.95, 1.7, PANEL, line=ACCENT2)
add_text(s, 7.15, 1.98, 5.4, 0.35, [[("⚙️  BACKEND — RENDER", 13, ACCENT2, True, False, True)]])
add_text(s, 7.15, 2.4, 5.4, 0.5, [[("https://empire-backend-pjvk.onrender.com", 14.5, INK, True, False, True)]])
add_text(s, 7.15, 2.95, 5.4, 0.4, [[("FastAPI · MediaPipe · sensors · AI · DB", 12, DIM, False, False, False)]])

add_rect(s, 0.55, 3.75, 12.25, 1.1, PANEL, line=GREEN)
add_text(s, 0.85, 3.95, 11.7, 0.35, [[("💾 SOURCE — GITHUB", 12.5, GREEN, True, False, True)]])
add_text(s, 0.85, 4.35, 11.7, 0.5, [[("github.com/YOGESH11K/visionlab-ai  ·  docs, firmware, tests, blueprints included", 14, INK, True, False, True)]])

add_text(s, 0.55, 5.25, 12.2, 1.6,
         [[("Engineered for production", 15, INK, True, False, False)],
          [("• 60+ automated tests · type-checked frontend · structured logging · graceful fallbacks", 13, DIM, False, False, False)],
          [("• Auto-deploy on git push · environment-driven config · secrets never committed", 13, DIM, False, False, False)],
          [("• Render free tier sleeps when idle (cold start ~30–60s) — upgrade or ping /api/health to keep warm", 13, ORANGE, False, False, False)]])

# =====================================================================
# SLIDE 15 — THANK YOU
# =====================================================================
s = prs.slides.add_slide(BLANK); set_bg(s)
for i in range(0, 14):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i*1.0), Inches(0), Pt(1), SH)
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(20, 26, 38); ln.line.fill.background(); ln.shadow.inherit = False
for j in range(0, 8):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(j*1.1), SW, Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(20, 26, 38); ln.line.fill.background(); ln.shadow.inherit = False
glow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(3.35), Inches(2.2), Pt(3))
glow.fill.solid(); glow.fill.fore_color.rgb = ACCENT; glow.line.fill.background(); glow.shadow.inherit = False

add_text(s, 0.55, 2.0, 12.2, 0.5, [[("THANK YOU", 52, INK, True, False, False)]])
add_text(s, 0.55, 3.0, 12.2, 0.5,
         [[("“A professional AI-powered electronics laboratory where computer vision understands the physical world, and you control hardware naturally.”", 16, DIM, False, True, False)]])
add_text(s, 0.55, 4.1, 12.2, 0.5, [[("👁 VisionLab AI  ·  🖐 Gestures  ·  🤖 AI  ·  🔌 Arduino/ESP32  ·  💡 Electronics  ·  📊 Data  ·  🧠 Learning  ·  🔧 Circuits", 15, ACCENT, True, False, False)]])
add_text(s, 0.55, 5.1, 12.2, 0.5, [[("Q & A  ·  Try it live:  https://visionlab-ai.vercel.app", 15, INK, True, False, False)]])

# =====================================================================
out = "C:/Users/Administrator/Documents/Python/robo/VisionLab_AI_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")